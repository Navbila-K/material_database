#!/usr/bin/env python3
"""
Create Professional Slide 1 for IIT Madras Funding Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_professional_slide1():
    """Create a stunning professional Slide 1"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Colors
    IIT_BLUE = RGBColor(0, 53, 102)
    ACCENT_ORANGE = RGBColor(255, 107, 53)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_BLUE = RGBColor(200, 220, 255)
    
    # Create blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # ===== BACKGROUND GRADIENT EFFECT =====
    # Dark blue background
    bg_main = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg_main.fill.solid()
    bg_main.fill.fore_color.rgb = IIT_BLUE
    bg_main.line.fill.background()
    
    # Decorative circles (modern design)
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-1), Inches(-1), Inches(4), Inches(4)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(0, 73, 122)  # Lighter blue
    circle1.fill.transparency = 0.5
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(7), Inches(5), Inches(4), Inches(4)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(255, 127, 73)  # Lighter orange
    circle2.fill.transparency = 0.5
    circle2.line.fill.background()
    
    # ===== ACCENT STRIPE AT BOTTOM =====
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.8), Inches(10), Inches(0.7)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_ORANGE
    accent_bar.line.fill.background()
    
    # ===== MAIN TITLE (LARGE & BOLD) =====
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    # Main title text
    p1 = title_frame.paragraphs[0]
    p1.text = "INTELLIGENT MATERIALS"
    p1.font.size = Pt(52)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = title_frame.add_paragraph()
    p2.text = "DATABASE SYSTEM"
    p2.font.size = Pt(52)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_ORANGE
    p2.alignment = PP_ALIGN.CENTER
    
    # ===== SUBTITLE (Tagline) =====
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    
    sub_p1 = subtitle_frame.paragraphs[0]
    sub_p1.text = "Advanced XML-Driven Material Property Management"
    sub_p1.font.size = Pt(20)
    sub_p1.font.color.rgb = LIGHT_BLUE
    sub_p1.alignment = PP_ALIGN.CENTER
    
    sub_p2 = subtitle_frame.add_paragraph()
    sub_p2.text = "with AI-Powered Visualization & Override Capabilities"
    sub_p2.font.size = Pt(20)
    sub_p2.font.color.rgb = LIGHT_BLUE
    sub_p2.alignment = PP_ALIGN.CENTER
    
    # ===== KEY HIGHLIGHT BOX =====
    highlight_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.5), Inches(5), Inches(5), Inches(0.6)
    )
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = RGBColor(50, 100, 150)  # Lighter blue
    highlight_box.line.color.rgb = ACCENT_ORANGE
    highlight_box.line.width = Pt(2)
    
    highlight_text = highlight_box.text_frame
    highlight_text.text = "🎯 Transforming Materials Research for India's Defense & Aerospace"
    highlight_text.paragraphs[0].font.size = Pt(16)
    highlight_text.paragraphs[0].font.bold = True
    highlight_text.paragraphs[0].font.color.rgb = WHITE
    highlight_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    highlight_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # ===== PRESENTED TO SECTION =====
    presented_box = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(6), Inches(0.8))
    presented_frame = presented_box.text_frame
    
    pres_p1 = presented_frame.paragraphs[0]
    pres_p1.text = "Presented to:"
    pres_p1.font.size = Pt(14)
    pres_p1.font.italic = True
    pres_p1.font.color.rgb = LIGHT_BLUE
    pres_p1.alignment = PP_ALIGN.CENTER
    
    pres_p2 = presented_frame.add_paragraph()
    pres_p2.text = "Indian Institute of Technology Madras"
    pres_p2.font.size = Pt(18)
    pres_p2.font.bold = True
    pres_p2.font.color.rgb = WHITE
    pres_p2.alignment = PP_ALIGN.CENTER
    
    pres_p3 = presented_frame.add_paragraph()
    pres_p3.text = "Research Funding Committee"
    pres_p3.font.size = Pt(16)
    pres_p3.font.color.rgb = LIGHT_BLUE
    pres_p3.alignment = PP_ALIGN.CENTER
    
    # ===== TEAM INFO (Bottom Left) =====
    team_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(3), Inches(0.4))
    team_frame = team_box.text_frame
    team_frame.text = "Presented by: [Your Team Name] | [Your Institution]"
    team_frame.paragraphs[0].font.size = Pt(11)
    team_frame.paragraphs[0].font.color.rgb = IIT_BLUE
    team_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # ===== DATE (Bottom Right) =====
    date_box = slide.shapes.add_textbox(Inches(7), Inches(7), Inches(2.5), Inches(0.4))
    date_frame = date_box.text_frame
    date_frame.text = "December 2025"
    date_frame.paragraphs[0].font.size = Pt(12)
    date_frame.paragraphs[0].font.bold = True
    date_frame.paragraphs[0].font.color.rgb = IIT_BLUE
    date_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    # ===== IIT LOGO PLACEHOLDER =====
    # Add a text box for logo instruction
    logo_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(0.3), Inches(1.2), Inches(0.5)
    )
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = WHITE
    logo_box.line.color.rgb = ACCENT_ORANGE
    logo_box.line.width = Pt(2)
    
    logo_text = logo_box.text_frame
    logo_text.text = "IIT-M\nLogo"
    logo_text.paragraphs[0].font.size = Pt(10)
    logo_text.paragraphs[0].font.bold = True
    logo_text.paragraphs[0].font.color.rgb = IIT_BLUE
    logo_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    logo_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Save to Desktop
    filename = '/Users/sridhars/Desktop/IIT_Madras_Professional_Slide1.pptx'
    prs.save(filename)
    print(f"✅ Professional Slide 1 created: {filename}")
    print(f"\n📝 Next steps:")
    print(f"   1. Open the file in PowerPoint/Keynote")
    print(f"   2. Replace '[Your Team Name]' with your actual team name")
    print(f"   3. Replace '[Your Institution]' with your college name")
    print(f"   4. Add IIT Madras logo in top-right placeholder")
    print(f"   5. Copy this slide to your existing presentation")
    
    return filename

if __name__ == "__main__":
    create_professional_slide1()
