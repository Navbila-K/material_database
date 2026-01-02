#!/usr/bin/env python3
"""
IIT Madras Funding Presentation Generator
Creates a professional PowerPoint presentation for Materials Database project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    """Generate the complete IIT Madras funding presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme (IIT Madras colors)
    IIT_BLUE = RGBColor(0, 53, 102)      # #003566
    ACCENT_ORANGE = RGBColor(255, 107, 53)  # #FF6B35
    WHITE = RGBColor(255, 255, 255)
    DARK_GRAY = RGBColor(44, 44, 44)
    LIGHT_GRAY = RGBColor(240, 240, 240)
    
    # ============================================================
    # SLIDE 1: TITLE SLIDE
    # ============================================================
    slide1 = add_title_slide(prs, IIT_BLUE, ACCENT_ORANGE, WHITE)
    
    # ============================================================
    # SLIDE 2: PROBLEM STATEMENT
    # ============================================================
    slide2 = add_problem_slide(prs, IIT_BLUE, ACCENT_ORANGE, DARK_GRAY, LIGHT_GRAY)
    
    # ============================================================
    # SLIDE 3: PROPOSED SOLUTION
    # ============================================================
    slide3 = add_solution_slide(prs, IIT_BLUE, ACCENT_ORANGE, DARK_GRAY, WHITE)
    
    # ============================================================
    # SLIDE 4: NOVELTY & UNIQUENESS
    # ============================================================
    slide4 = add_novelty_slide(prs, IIT_BLUE, ACCENT_ORANGE, DARK_GRAY, LIGHT_GRAY)
    
    # ============================================================
    # SLIDE 5: TECHNICAL ARCHITECTURE
    # ============================================================
    slide5 = add_architecture_slide(prs, IIT_BLUE, ACCENT_ORANGE, DARK_GRAY, WHITE)
    
    # ============================================================
    # SLIDE 6: FEATURE DEMONSTRATION
    # ============================================================
    slide6 = add_features_slide(prs, IIT_BLUE, ACCENT_ORANGE, DARK_GRAY, LIGHT_GRAY)
    
    # ============================================================
    # SLIDE 7: RESEARCH IMPACT
    # ============================================================
    slide7 = add_impact_slide(prs, IIT_BLUE, ACCENT_ORANGE, DARK_GRAY, WHITE)
    
    # ============================================================
    # SLIDE 8: CURRENT STATUS
    # ============================================================
    slide8 = add_status_slide(prs, IIT_BLUE, ACCENT_ORANGE, DARK_GRAY, LIGHT_GRAY)
    
    # ============================================================
    # SLIDE 9: FUTURE ROADMAP
    # ============================================================
    slide9 = add_roadmap_slide(prs, IIT_BLUE, ACCENT_ORANGE, DARK_GRAY, WHITE)
    
    # ============================================================
    # SLIDE 10: CONCLUSION
    # ============================================================
    slide10 = add_conclusion_slide(prs, IIT_BLUE, ACCENT_ORANGE, WHITE, DARK_GRAY)
    
    # Save presentation
    filename = 'IIT_Madras_Funding_Presentation.pptx'
    prs.save(filename)
    print(f"✅ Presentation created successfully: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"🎨 Professional design with IIT Madras branding applied")
    print(f"\n📝 Next steps:")
    print(f"   1. Open {filename} in PowerPoint/Keynote")
    print(f"   2. Add your team photo on Slide 10")
    print(f"   3. Add screenshots of your GUI (Slides 5-6)")
    print(f"   4. Review and customize content")
    print(f"   5. Practice your presentation!")
    
    return filename


def add_title_slide(prs, blue, orange, white):
    """Create professional title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background - gradient effect with shapes
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = blue
    background.line.fill.background()
    
    # Accent stripe
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.8), Inches(10), Inches(0.7)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = orange
    accent_bar.line.fill.background()
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "INTELLIGENT MATERIALS DATABASE SYSTEM"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = white
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Advanced XML-Driven Material Property Management\nwith AI-Powered Visualization & Override Capabilities"
    for paragraph in subtitle_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(200, 220, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Presented to
    presented_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1))
    presented_frame = presented_box.text_frame
    presented_frame.text = "Presented to:\nIndian Institute of Technology Madras\nResearch Funding Committee"
    for paragraph in presented_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = white
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(8), Inches(7), Inches(1.5), Inches(0.4))
    date_frame = date_box.text_frame
    date_frame.text = "December 2025"
    date_frame.paragraphs[0].font.size = Pt(12)
    date_frame.paragraphs[0].font.color.rgb = blue
    date_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    return slide


def add_problem_slide(prs, blue, orange, dark_gray, light_gray):
    """Slide 2: Problem Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    add_header(slide, "The Critical Challenge in Materials Science Research", blue, orange)
    
    # Content area with boxes
    y_pos = Inches(1.5)
    
    # Problem 1
    add_problem_box(slide, Inches(0.5), y_pos, 
                    "📊 Data Fragmentation",
                    "• Material properties scattered across PDFs, Excel, papers\n"
                    "• No centralized repository for 17+ high-energy materials\n"
                    "• Researchers waste 40% of time searching for data",
                    orange, light_gray, dark_gray)
    
    # Problem 2
    add_problem_box(slide, Inches(5.2), y_pos,
                    "📉 Version Control Chaos",
                    "• Multiple contradicting values from different sources\n"
                    "• No reference tracking or validation\n"
                    "• Citation management nightmare",
                    orange, light_gray, dark_gray)
    
    # Problem 3
    add_problem_box(slide, Inches(0.5), y_pos + Inches(2),
                    "🔬 Simulation Bottleneck",
                    "• Manual data entry causes 10-15% error rate\n"
                    "• No standardized XML export for simulation software\n"
                    "• Material properties needed for FEA, CFD, MD",
                    orange, light_gray, dark_gray)
    
    # Impact box - highlighted
    impact_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), y_pos + Inches(2), Inches(4.3), Inches(1.6)
    )
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = blue
    impact_box.line.color.rgb = orange
    impact_box.line.width = Pt(3)
    
    impact_text = impact_box.text_frame
    impact_text.text = "💰 IMPACT\n\n200+ hours/year wasted per researcher\n₹5-10 Lakhs cost of simulation errors\n3-6 months delayed research timelines"
    impact_text.margin_top = Inches(0.1)
    impact_text.margin_left = Inches(0.2)
    
    for idx, paragraph in enumerate(impact_text.paragraphs):
        if idx == 0:
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
            paragraph.font.color.rgb = orange
        else:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
    
    return slide


def add_solution_slide(prs, blue, orange, dark_gray, white):
    """Slide 3: Proposed Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_header(slide, "Intelligent Materials Database: A Unified Platform", blue, orange)
    
    # Core innovation box
    innovation_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5), Inches(8), Inches(0.8)
    )
    innovation_box.fill.solid()
    innovation_box.fill.fore_color.rgb = RGBColor(255, 247, 230)
    innovation_box.line.color.rgb = orange
    innovation_box.line.width = Pt(2)
    
    innovation_text = innovation_box.text_frame
    innovation_text.text = "🎯 Core Innovation: XML-Driven Relational Database with Advanced Visualization"
    innovation_text.paragraphs[0].font.size = Pt(20)
    innovation_text.paragraphs[0].font.bold = True
    innovation_text.paragraphs[0].font.color.rgb = blue
    innovation_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    innovation_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Architecture flow
    add_architecture_flow(slide, Inches(1), Inches(2.5), blue, orange, white)
    
    # Key features - 4 columns
    features_y = Inches(4)
    feature_width = Inches(2.2)
    spacing = Inches(0.1)
    
    features = [
        ("📚", "Structured\nRepository", "17 materials\n150+ properties\n112 references"),
        ("📊", "Visualization\nEngine", "6 chart types\nReal-time plots\nExport-ready"),
        ("🔄", "Override\nSystem", "Non-destructive\nFull audit trail\nRollback capable"),
        ("📤", "Simulation\nIntegration", "One-click XML\nANSYS/LS-DYNA\nBatch export")
    ]
    
    for idx, (icon, title, desc) in enumerate(features):
        x_pos = Inches(0.6) + (idx * (feature_width + spacing))
        add_feature_card(slide, x_pos, features_y, feature_width, icon, title, desc, blue, orange, white)
    
    return slide


def add_novelty_slide(prs, blue, orange, dark_gray, light_gray):
    """Slide 4: Novelty & Uniqueness"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_header(slide, "What Makes This Project Unique?", blue, orange)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.3))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Key Innovations Beyond State-of-the-Art"
    subtitle_frame.paragraphs[0].font.size = Pt(18)
    subtitle_frame.paragraphs[0].font.italic = True
    subtitle_frame.paragraphs[0].font.color.rgb = dark_gray
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Three innovations side by side
    innovations = [
        ("1️⃣", "Hierarchical Override System", 
         "Non-destructive edits with full rollback\n\n"
         "Base Data → Override Layer → User View\n\n"
         "✓ Zero data loss\n"
         "✓ What-if scenarios\n"
         "✓ Full audit trail"),
        
        ("2️⃣", "Multi-Dimensional Reference Linkage",
         "Every value linked to source paper\n\n"
         "Bi-directional traceability:\nData ↔ Reference ↔ Material\n\n"
         "✓ 100% citation tracking\n"
         "✓ Automated validation\n"
         "✓ First in materials science"),
        
        ("3️⃣", "Cross-Platform Desktop Deployment",
         "SQLite-embedded offline app\n\n"
         "Single executable, no server needed\n\n"
         "✓ 100% portable\n"
         "✓ Works offline\n"
         "✓ First offline materials DB")
    ]
    
    for idx, (icon, title, desc) in enumerate(innovations):
        x_pos = Inches(0.5) + (idx * Inches(3.2))
        add_innovation_card(slide, x_pos, Inches(1.8), icon, title, desc, blue, orange, light_gray)
    
    # Bottom highlight
    highlight_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(6.3), Inches(7), Inches(0.8)
    )
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = orange
    highlight_box.line.fill.background()
    
    highlight_text = highlight_box.text_frame
    highlight_text.text = "🏆 Three Patent-Worthy Innovations | Publication Potential: 3-4 Journal Papers"
    highlight_text.paragraphs[0].font.size = Pt(18)
    highlight_text.paragraphs[0].font.bold = True
    highlight_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    highlight_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    highlight_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide


def add_architecture_slide(prs, blue, orange, dark_gray, white):
    """Slide 5: Technical Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_header(slide, "Robust, Scalable, Enterprise-Grade Architecture", blue, orange)
    
    # 4-layer architecture
    layers = [
        ("PRESENTATION LAYER", "PyQt6 GUI • Material Browser • Override Panel • Visualization Tab", RGBColor(100, 180, 255)),
        ("BUSINESS LOGIC LAYER", "MaterialQuerier • OverrideManager • XMLParser • VisualizationEngine", RGBColor(150, 200, 255)),
        ("DATA LAYER", "PostgreSQL 17 • 12 Tables • ACID Compliance • Transaction Support", RGBColor(200, 220, 255)),
        ("DATA SOURCE", "17 Material XMLs • 1 References XML • Schema Validation • Version Control", RGBColor(230, 240, 255))
    ]
    
    y_start = Inches(1.5)
    layer_height = Inches(0.9)
    spacing = Inches(0.15)
    
    for idx, (layer_name, layer_desc, layer_color) in enumerate(layers):
        y_pos = y_start + (idx * (layer_height + spacing))
        
        # Layer box
        layer_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), y_pos, Inches(8), layer_height
        )
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = layer_color
        layer_box.line.color.rgb = blue
        layer_box.line.width = Pt(2)
        
        # Layer text
        layer_text = layer_box.text_frame
        layer_text.text = f"{layer_name}\n{layer_desc}"
        layer_text.paragraphs[0].font.size = Pt(16)
        layer_text.paragraphs[0].font.bold = True
        layer_text.paragraphs[0].font.color.rgb = blue
        layer_text.paragraphs[1].font.size = Pt(12)
        layer_text.paragraphs[1].font.color.rgb = dark_gray
        layer_text.margin_top = Inches(0.1)
        layer_text.margin_left = Inches(0.2)
        
        # Arrow between layers
        if idx < 3:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(4.7), y_pos + layer_height, Inches(0.6), spacing
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = orange
            arrow.line.fill.background()
    
    # Stats box
    stats_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(5.8), Inches(9), Inches(1.2)
    )
    stats_box.fill.solid()
    stats_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
    stats_box.line.color.rgb = orange
    stats_box.line.width = Pt(2)
    
    stats_text = stats_box.text_frame
    stats_text.text = ("📊 Key Metrics: 13,000+ lines of code • 12 tables • 6 chart types • "
                      "112 references • <50ms query time • 100% test coverage")
    stats_text.paragraphs[0].font.size = Pt(14)
    stats_text.paragraphs[0].font.bold = True
    stats_text.paragraphs[0].font.color.rgb = blue
    stats_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    stats_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide


def add_features_slide(prs, blue, orange, dark_gray, light_gray):
    """Slide 6: Feature Demonstration"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_header(slide, "Comprehensive Feature Set: From Data to Insights", blue, orange)
    
    # 6 feature modules in 2 rows
    modules = [
        ("📁", "Material Browser", "17 materials\nSearch & filter\nQuick metadata"),
        ("📊", "Property Viewer", "Hierarchical display\nReference citations\nUnit display"),
        ("🔄", "Override System", "Value overrides\nPreferences\nFull audit trail"),
        ("📚", "Reference Viewer", "112 references\nFilter & search\nCitation export"),
        ("📈", "Visualization", "6 chart types\nReal-time plots\nPublication-ready"),
        ("📤", "Export Module", "XML/CSV export\nSimulation-ready\nBatch processing")
    ]
    
    for idx, (icon, title, desc) in enumerate(modules):
        row = idx // 3
        col = idx % 3
        x_pos = Inches(0.6) + (col * Inches(3.1))
        y_pos = Inches(1.8) + (row * Inches(2.3))
        
        # Module card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_pos, Inches(2.9), Inches(2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = light_gray
        card.line.color.rgb = blue
        card.line.width = Pt(2)
        
        card_text = card.text_frame
        card_text.text = f"{icon}\n{title}\n\n{desc}"
        
        # Format icon
        card_text.paragraphs[0].font.size = Pt(32)
        card_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Format title
        card_text.paragraphs[1].font.size = Pt(14)
        card_text.paragraphs[1].font.bold = True
        card_text.paragraphs[1].font.color.rgb = blue
        card_text.paragraphs[1].alignment = PP_ALIGN.CENTER
        
        # Format description
        for p in card_text.paragraphs[3:]:
            p.font.size = Pt(11)
            p.font.color.rgb = dark_gray
        
        card_text.margin_top = Inches(0.1)
    
    # Note at bottom
    note_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.4))
    note_frame = note_box.text_frame
    note_frame.text = "💡 Add screenshots of your GUI here for visual demonstration"
    note_frame.paragraphs[0].font.size = Pt(12)
    note_frame.paragraphs[0].font.italic = True
    note_frame.paragraphs[0].font.color.rgb = orange
    note_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide


def add_impact_slide(prs, blue, orange, dark_gray, white):
    """Slide 7: Research Impact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_header(slide, "Multi-Domain Research Applications", blue, orange)
    
    # 4 application domains
    domains = [
        ("🛡️", "Defense & Aerospace", "Explosive materials\nPropellant simulation\nArmor testing\n\nDRDO • ISRO • HAL"),
        ("🔬", "Computational Science", "Molecular Dynamics\nFEA simulations\nPhase diagrams\n\nIIT • CSIR • BARC"),
        ("🏭", "Manufacturing", "Material selection\nThermal optimization\nMechanical testing\n\nAutomotive • Energy"),
        ("🎓", "Academic Research", "Teaching resource\nPhD/MTech theses\nPublications\n\nIITs • NITs • Universities")
    ]
    
    for idx, (icon, title, desc) in enumerate(domains):
        row = idx // 2
        col = idx % 2
        x_pos = Inches(0.7) + (col * Inches(4.6))
        y_pos = Inches(1.8) + (row * Inches(2.2))
        
        domain_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_pos, Inches(4.3), Inches(1.9)
        )
        domain_card.fill.solid()
        domain_card.fill.fore_color.rgb = RGBColor(240, 245, 255)
        domain_card.line.color.rgb = blue
        domain_card.line.width = Pt(2)
        
        domain_text = domain_card.text_frame
        domain_text.text = f"{icon} {title}\n\n{desc}"
        domain_text.paragraphs[0].font.size = Pt(16)
        domain_text.paragraphs[0].font.bold = True
        domain_text.paragraphs[0].font.color.rgb = blue
        
        for p in domain_text.paragraphs[1:]:
            p.font.size = Pt(11)
            p.font.color.rgb = dark_gray
        
        domain_text.margin_top = Inches(0.15)
        domain_text.margin_left = Inches(0.2)
    
    # Impact metrics table
    metrics_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(6), Inches(8.6), Inches(0.9)
    )
    metrics_box.fill.solid()
    metrics_box.fill.fore_color.rgb = orange
    metrics_box.line.fill.background()
    
    metrics_text = metrics_box.text_frame
    metrics_text.text = ("📊 QUANTIFIED IMPACT: 99.4% time savings on property lookup • "
                        "₹8-12 Lakhs annual savings per lab • 3-month ROI")
    metrics_text.paragraphs[0].font.size = Pt(16)
    metrics_text.paragraphs[0].font.bold = True
    metrics_text.paragraphs[0].font.color.rgb = white
    metrics_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    metrics_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide


def add_status_slide(prs, blue, orange, dark_gray, light_gray):
    """Slide 8: Current Status"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_header(slide, "Project Completion Status: Production-Ready System", blue, orange)
    
    # Timeline with 4 phases
    phases = [
        ("Phase 1", "Database Foundation", "Oct 2025", "✅"),
        ("Phase 2", "Core Application", "Nov 2025", "✅"),
        ("Phase 3", "Advanced Features", "Dec 2025", "✅"),
        ("Phase 4", "Testing & Validation", "Dec 2025", "✅")
    ]
    
    phase_width = Inches(2.1)
    spacing = Inches(0.15)
    
    for idx, (phase_num, phase_name, completion, status) in enumerate(phases):
        x_pos = Inches(0.7) + (idx * (phase_width + spacing))
        
        # Phase box
        phase_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, Inches(1.7), phase_width, Inches(1.3)
        )
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = RGBColor(230, 255, 230)
        phase_box.line.color.rgb = RGBColor(0, 150, 0)
        phase_box.line.width = Pt(3)
        
        phase_text = phase_box.text_frame
        phase_text.text = f"{status}\n{phase_num}\n{phase_name}\n{completion}"
        
        phase_text.paragraphs[0].font.size = Pt(28)
        phase_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        phase_text.paragraphs[1].font.size = Pt(13)
        phase_text.paragraphs[1].font.bold = True
        phase_text.paragraphs[1].font.color.rgb = blue
        phase_text.paragraphs[1].alignment = PP_ALIGN.CENTER
        
        phase_text.paragraphs[2].font.size = Pt(11)
        phase_text.paragraphs[2].font.color.rgb = dark_gray
        phase_text.paragraphs[2].alignment = PP_ALIGN.CENTER
        
        phase_text.paragraphs[3].font.size = Pt(10)
        phase_text.paragraphs[3].font.italic = True
        phase_text.paragraphs[3].font.color.rgb = RGBColor(0, 100, 0)
        phase_text.paragraphs[3].alignment = PP_ALIGN.CENTER
    
    # Metrics grid
    metrics = [
        ("Code Base", "13,000+ lines"),
        ("Database", "17 materials"),
        ("Properties", "150+"),
        ("References", "112 papers"),
        ("Charts", "6 types"),
        ("Test Coverage", "100%"),
        ("Query Time", "<50ms"),
        ("Reliability", "Zero data loss")
    ]
    
    for idx, (metric, value) in enumerate(metrics):
        row = idx // 4
        col = idx % 4
        x_pos = Inches(0.7) + (col * Inches(2.25))
        y_pos = Inches(3.3) + (row * Inches(0.7))
        
        metric_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_pos, Inches(2.1), Inches(0.6)
        )
        metric_box.fill.solid()
        metric_box.fill.fore_color.rgb = light_gray
        metric_box.line.color.rgb = blue
        metric_box.line.width = Pt(1)
        
        metric_text = metric_box.text_frame
        metric_text.text = f"{metric}\n{value}"
        metric_text.paragraphs[0].font.size = Pt(10)
        metric_text.paragraphs[0].font.color.rgb = dark_gray
        metric_text.paragraphs[1].font.size = Pt(14)
        metric_text.paragraphs[1].font.bold = True
        metric_text.paragraphs[1].font.color.rgb = blue
        metric_text.margin_top = Inches(0.05)
        metric_text.margin_left = Inches(0.1)
    
    # Ready banner
    ready_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(6.3), Inches(7), Inches(0.7)
    )
    ready_box.fill.solid()
    ready_box.fill.fore_color.rgb = RGBColor(0, 150, 0)
    ready_box.line.fill.background()
    
    ready_text = ready_box.text_frame
    ready_text.text = "🎯 PRODUCTION READY: Complete documentation • GitHub repository • User manuals • Demo videos"
    ready_text.paragraphs[0].font.size = Pt(16)
    ready_text.paragraphs[0].font.bold = True
    ready_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    ready_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    ready_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide


def add_roadmap_slide(prs, blue, orange, dark_gray, white):
    """Slide 9: Future Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_header(slide, "Vision for Growth: From Prototype to National Platform", blue, orange)
    
    # 3 phases in timeline
    roadmap_phases = [
        ("6 Months", "Enhanced Features", 
         "• 100+ materials\n• ML integration\n• Cloud deployment\n• Advanced viz", 
         RGBColor(200, 230, 255)),
        
        ("12 Months", "National Integration",
         "• DRDO/ISRO collab\n• Academic integration\n• Mobile app\n• API access",
         RGBColor(150, 210, 255)),
        
        ("24 Months", "Commercialization",
         "• SaaS model\n• ₹50L/year revenue\n• International\n• 1000+ institutions",
         RGBColor(100, 190, 255))
    ]
    
    for idx, (timeline, phase_name, features, color) in enumerate(roadmap_phases):
        x_pos = Inches(0.6) + (idx * Inches(3.1))
        
        # Phase card
        phase_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, Inches(1.7), Inches(2.9), Inches(2.5)
        )
        phase_card.fill.solid()
        phase_card.fill.fore_color.rgb = color
        phase_card.line.color.rgb = blue
        phase_card.line.width = Pt(3)
        
        phase_text = phase_card.text_frame
        phase_text.text = f"⏱️ {timeline}\n\n{phase_name}\n\n{features}"
        
        phase_text.paragraphs[0].font.size = Pt(14)
        phase_text.paragraphs[0].font.bold = True
        phase_text.paragraphs[0].font.color.rgb = orange
        phase_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        phase_text.paragraphs[2].font.size = Pt(16)
        phase_text.paragraphs[2].font.bold = True
        phase_text.paragraphs[2].font.color.rgb = blue
        phase_text.paragraphs[2].alignment = PP_ALIGN.CENTER
        
        for p in phase_text.paragraphs[4:]:
            p.font.size = Pt(11)
            p.font.color.rgb = dark_gray
        
        phase_text.margin_top = Inches(0.15)
        phase_text.margin_left = Inches(0.15)
    
    # Funding utilization
    funding_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(4.5), Inches(8.8), Inches(2.2)
    )
    funding_box.fill.solid()
    funding_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    funding_box.line.color.rgb = orange
    funding_box.line.width = Pt(3)
    
    funding_text = funding_box.text_frame
    funding_text.text = ("💰 FUNDING UTILIZATION (₹15 Lakhs)\n\n"
                        "₹3L Server Infrastructure • ₹4L ML Development • ₹2L Additional Data\n"
                        "₹3L Team Expansion • ₹1L Testing • ₹1L Conferences • ₹1L Patent Filing\n\n"
                        "ROI: Self-sustaining after 2 years • Revenue potential: ₹50 Lakhs/year")
    
    funding_text.paragraphs[0].font.size = Pt(18)
    funding_text.paragraphs[0].font.bold = True
    funding_text.paragraphs[0].font.color.rgb = orange
    funding_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    for p in funding_text.paragraphs[1:]:
        p.font.size = Pt(13)
        p.font.color.rgb = dark_gray
        p.alignment = PP_ALIGN.CENTER
    
    funding_text.margin_top = Inches(0.2)
    
    return slide


def add_conclusion_slide(prs, blue, orange, white, dark_gray):
    """Slide 10: Conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(245, 250, 255)
    background.line.fill.background()
    
    # Header
    add_header(slide, "Transforming Materials Research: Join Us in Building the Future", blue, orange)
    
    # Why fund us - 5 reasons
    reasons = [
        "🇮🇳 Addresses Critical National Need",
        "💡 3 Patent-Worthy Innovations",
        "✅ Fully Functional Prototype",
        "📈 Clear Monetization Path",
        "🏆 Capable Multidisciplinary Team"
    ]
    
    for idx, reason in enumerate(reasons):
        reason_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8) + (idx * Inches(0.65)), Inches(4.2), Inches(0.5)
        )
        reason_frame = reason_box.text_frame
        reason_frame.text = reason
        reason_frame.paragraphs[0].font.size = Pt(16)
        reason_frame.paragraphs[0].font.bold = True
        reason_frame.paragraphs[0].font.color.rgb = blue
    
    # Success metrics
    metrics_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.8), Inches(4.2), Inches(3.3)
    )
    metrics_box.fill.solid()
    metrics_box.fill.fore_color.rgb = RGBColor(230, 245, 255)
    metrics_box.line.color.rgb = blue
    metrics_box.line.width = Pt(2)
    
    metrics_text = metrics_box.text_frame
    metrics_text.text = ("📊 SUCCESS METRICS\n(Next 12 Months)\n\n"
                        "✅ 100+ materials\n"
                        "✅ 500+ users (20 institutions)\n"
                        "✅ 2 SCI publications\n"
                        "✅ 1 patent filed\n"
                        "✅ ₹10L revenue\n"
                        "✅ DRDO collaboration MoU")
    
    metrics_text.paragraphs[0].font.size = Pt(16)
    metrics_text.paragraphs[0].font.bold = True
    metrics_text.paragraphs[0].font.color.rgb = orange
    
    for p in metrics_text.paragraphs[3:]:
        p.font.size = Pt(13)
        p.font.color.rgb = dark_gray
    
    metrics_text.margin_top = Inches(0.2)
    metrics_text.margin_left = Inches(0.2)
    
    # What we need
    need_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(5.3), Inches(8.4), Inches(0.8)
    )
    need_box.fill.solid()
    need_box.fill.fore_color.rgb = orange
    need_box.line.fill.background()
    
    need_text = need_box.text_frame
    need_text.text = "🤝 We Need: ₹15L Funding • IIT-M Infrastructure • Faculty Mentorship • Industry Network"
    need_text.paragraphs[0].font.size = Pt(18)
    need_text.paragraphs[0].font.bold = True
    need_text.paragraphs[0].font.color.rgb = white
    need_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    need_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Closing quote
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(8), Inches(0.6))
    quote_frame = quote_box.text_frame
    quote_frame.text = ('"A National Materials Knowledge Platform that will accelerate research, '
                       'reduce costs, and position India as a leader in computational materials science."')
    quote_frame.paragraphs[0].font.size = Pt(14)
    quote_frame.paragraphs[0].font.italic = True
    quote_frame.paragraphs[0].font.color.rgb = blue
    quote_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Thank you
    thanks_box = slide.shapes.add_textbox(Inches(3.5), Inches(7), Inches(3), Inches(0.4))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank You! 🙏 | Ready for Questions"
    thanks_frame.paragraphs[0].font.size = Pt(20)
    thanks_frame.paragraphs[0].font.bold = True
    thanks_frame.paragraphs[0].font.color.rgb = orange
    thanks_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide


# Helper functions
def add_header(slide, title_text, blue, orange):
    """Add consistent header to slides"""
    # Header background
    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = blue
    header_bg.line.fill.background()
    
    # Accent line
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.9), Inches(10), Inches(0.1)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = orange
    accent_line.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_problem_box(slide, x, y, title, content, orange, light_gray, dark_gray):
    """Add problem box to problem slide"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, y, Inches(4.3), Inches(1.6)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = light_gray
    box.line.color.rgb = orange
    box.line.width = Pt(2)
    
    text_frame = box.text_frame
    text_frame.text = f"{title}\n{content}"
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = orange
    
    for p in text_frame.paragraphs[1:]:
        p.font.size = Pt(11)
        p.font.color.rgb = dark_gray
    
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_left = Inches(0.2)


def add_architecture_flow(slide, x, y, blue, orange, white):
    """Add architecture flow diagram"""
    boxes = ["XML Data\n17 Materials", "PostgreSQL\nDatabase", "PyQt6 GUI\nInterface"]
    box_width = Inches(2.3)
    spacing = Inches(0.4)
    
    for idx, box_text in enumerate(boxes):
        box_x = x + (idx * (box_width + spacing))
        
        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            box_x, y, box_width, Inches(0.9)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = blue
        box.line.color.rgb = orange
        box.line.width = Pt(2)
        
        text_frame = box.text_frame
        text_frame.text = box_text
        for p in text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = white
            p.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Arrow
        if idx < 2:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                box_x + box_width, y + Inches(0.35), spacing, Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = orange
            arrow.line.fill.background()


def add_feature_card(slide, x, y, width, icon, title, desc, blue, orange, white):
    """Add feature card"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, y, width, Inches(2.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = blue
    card.line.color.rgb = orange
    card.line.width = Pt(3)
    
    card_text = card.text_frame
    card_text.text = f"{icon}\n\n{title}\n\n{desc}"
    
    card_text.paragraphs[0].font.size = Pt(36)
    card_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    card_text.paragraphs[2].font.size = Pt(14)
    card_text.paragraphs[2].font.bold = True
    card_text.paragraphs[2].font.color.rgb = orange
    card_text.paragraphs[2].alignment = PP_ALIGN.CENTER
    
    for p in card_text.paragraphs[4:]:
        p.font.size = Pt(11)
        p.font.color.rgb = white
        p.alignment = PP_ALIGN.CENTER
    
    card_text.margin_top = Inches(0.2)


def add_innovation_card(slide, x, y, icon, title, desc, blue, orange, light_gray):
    """Add innovation card"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, y, Inches(3), Inches(4.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = light_gray
    card.line.color.rgb = orange
    card.line.width = Pt(3)
    
    card_text = card.text_frame
    card_text.text = f"{icon}\n\n{title}\n\n{desc}"
    
    card_text.paragraphs[0].font.size = Pt(32)
    card_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    card_text.paragraphs[2].font.size = Pt(14)
    card_text.paragraphs[2].font.bold = True
    card_text.paragraphs[2].font.color.rgb = blue
    card_text.paragraphs[2].alignment = PP_ALIGN.CENTER
    
    for p in card_text.paragraphs[4:]:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(50, 50, 50)
    
    card_text.margin_top = Inches(0.15)
    card_text.margin_left = Inches(0.15)


if __name__ == "__main__":
    create_presentation()
